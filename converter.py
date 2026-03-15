"""Convert ECharts option dicts to native, editable PowerPoint charts."""

from io import BytesIO
from typing import List, Optional, Set, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

SUPPORTED_TYPES: Set[str] = {
    "bar",
    "line",
    "pie",
    "scatter",
    "effectScatter",
    "radar",
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

CHART_LEFT = Inches(0.5)
CHART_TOP = Inches(1.2)
CHART_WIDTH = Inches(12.333)
CHART_HEIGHT = Inches(5.8)


def is_exportable(option: dict) -> bool:
    """Return True if every series type in *option* can be exported to PPTX."""
    series_list = option.get("series", [])
    return bool(series_list) and all(
        s.get("type", "bar") in SUPPORTED_TYPES for s in series_list
    )


def _extract_title(option: dict) -> Optional[str]:
    title = option.get("title")
    if isinstance(title, dict):
        return title.get("text")
    return None


def _extract_categories(option: dict, series_list: List[dict]) -> List[str]:
    """Extract category labels from xAxis.data or from series data point names."""
    x_axis = option.get("xAxis")
    if isinstance(x_axis, dict) and "data" in x_axis:
        return [str(c) for c in x_axis["data"]]
    if isinstance(x_axis, list) and x_axis and "data" in x_axis[0]:
        return [str(c) for c in x_axis[0]["data"]]

    for s in series_list:
        data = s.get("data", [])
        if data and isinstance(data[0], dict) and "name" in data[0]:
            return [str(d["name"]) for d in data]

    max_len = max((len(s.get("data", [])) for s in series_list), default=0)
    return [str(i + 1) for i in range(max_len)]


def _extract_series_values(series: dict) -> List[float]:
    """Extract numeric values from a series data array."""
    values = []
    for item in series.get("data", []):
        if isinstance(item, dict):
            values.append(float(item.get("value", 0)))
        elif isinstance(item, (int, float)):
            values.append(float(item))
        else:
            values.append(0.0)
    return values


def _extract_scatter_points(series: dict) -> List[Tuple[float, float]]:
    """Extract (x, y) pairs from scatter series data."""
    points = []
    for item in series.get("data", []):
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            points.append((float(item[0]), float(item[1])))
    return points


def _add_title_textbox(slide, text: str):
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8)
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True


def _is_horizontal_bar(option: dict) -> bool:
    """Horizontal bar when yAxis carries the categories."""
    y_axis = option.get("yAxis")
    if isinstance(y_axis, dict):
        return y_axis.get("type") == "category"
    if isinstance(y_axis, list) and y_axis:
        return y_axis[0].get("type") == "category"
    return False


def _has_stack(series_list: List[dict]) -> bool:
    return any(s.get("stack") for s in series_list)


def _has_area_style(series_list: List[dict]) -> bool:
    return any("areaStyle" in s for s in series_list)


def _is_donut(option: dict, series_list: List[dict]) -> bool:
    """A donut is a pie whose first radius value is non-zero."""
    series = series_list[0]
    radius = series.get("radius")
    if isinstance(radius, (list, tuple)) and len(radius) >= 1:
        inner = radius[0]
        if isinstance(inner, str):
            return inner.replace("%", "").strip() not in ("", "0")
        return float(inner) > 0
    return False


def _resolve_category_chart_type(
    chart_type_str: str,
    option: dict,
    series_list: List[dict],
) -> "XL_CHART_TYPE":
    if chart_type_str == "bar":
        if _is_horizontal_bar(option):
            if _has_stack(series_list):
                return XL_CHART_TYPE.BAR_STACKED
            return XL_CHART_TYPE.BAR_CLUSTERED
        if _has_stack(series_list):
            return XL_CHART_TYPE.COLUMN_STACKED
        return XL_CHART_TYPE.COLUMN_CLUSTERED

    if chart_type_str == "line":
        if _has_area_style(series_list):
            if _has_stack(series_list):
                return XL_CHART_TYPE.AREA_STACKED
            return XL_CHART_TYPE.AREA
        return XL_CHART_TYPE.LINE_MARKERS

    return XL_CHART_TYPE.COLUMN_CLUSTERED


def _build_category_chart(
    prs: Presentation,
    option: dict,
    series_list: List[dict],
    chart_type_str: str,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _extract_title(option)
    if title:
        _add_title_textbox(slide, title)

    if _is_horizontal_bar(option):
        categories = _extract_categories(
            {"xAxis": option.get("yAxis"), "yAxis": option.get("xAxis")},
            series_list,
        )
    else:
        categories = _extract_categories(option, series_list)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", "Series"), _extract_series_values(s))

    pptx_type = _resolve_category_chart_type(chart_type_str, option, series_list)
    graphic_frame = slide.shapes.add_chart(
        pptx_type, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, chart_data
    )
    chart = graphic_frame.chart

    if len(series_list) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


def _build_pie_chart(
    prs: Presentation,
    option: dict,
    series_list: List[dict],
    as_donut: bool = False,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _extract_title(option)
    if title:
        _add_title_textbox(slide, title)

    series = series_list[0]
    categories = _extract_categories(option, [series])
    values = _extract_series_values(series)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series.get("name", "Series"), values)

    pptx_type = XL_CHART_TYPE.DOUGHNUT if as_donut else XL_CHART_TYPE.PIE
    graphic_frame = slide.shapes.add_chart(
        pptx_type, CHART_LEFT, CHART_TOP, CHART_WIDTH, CHART_HEIGHT, chart_data
    )
    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


def _build_scatter_chart(prs: Presentation, option: dict, series_list: List[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _extract_title(option)
    if title:
        _add_title_textbox(slide, title)

    chart_data = XyChartData()
    for s in series_list:
        xy_series = chart_data.add_series(s.get("name", "Series"))
        for x, y in _extract_scatter_points(s):
            xy_series.add_data_point(x, y)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        CHART_LEFT,
        CHART_TOP,
        CHART_WIDTH,
        CHART_HEIGHT,
        chart_data,
    )
    chart = graphic_frame.chart

    if len(series_list) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


def _build_radar_chart(prs: Presentation, option: dict, series_list: List[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _extract_title(option)
    if title:
        _add_title_textbox(slide, title)

    indicators = option.get("radar", {})
    if isinstance(indicators, list):
        indicators = indicators[0] if indicators else {}
    indicator_list = indicators.get("indicator", [])
    categories = [str(ind.get("name", "")) for ind in indicator_list]

    chart_data = CategoryChartData()
    chart_data.categories = categories

    radar_series = series_list[0]
    for entry in radar_series.get("data", []):
        name = entry.get("name", "Series")
        values = [float(v) for v in entry.get("value", [])]
        chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_MARKERS,
        CHART_LEFT,
        CHART_TOP,
        CHART_WIDTH,
        CHART_HEIGHT,
        chart_data,
    )
    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False


def echarts_to_pptx(option: dict) -> bytes:
    """Convert an ECharts option dict to a PowerPoint file returned as bytes.

    The resulting PPTX contains native, editable charts backed by embedded
    Excel worksheets.  Supported ECharts series types: bar, line, pie,
    scatter, effectScatter, and radar.  Sub-variants (donut, area, stacked,
    horizontal bar) are auto-detected from series properties.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    series_list = option.get("series", [])
    if not series_list:
        raise ValueError("ECharts option must contain at least one series")

    chart_type = series_list[0].get("type", "bar")

    if chart_type == "pie":
        _build_pie_chart(prs, option, series_list, as_donut=_is_donut(option, series_list))
    elif chart_type in ("scatter", "effectScatter"):
        _build_scatter_chart(prs, option, series_list)
    elif chart_type == "radar":
        _build_radar_chart(prs, option, series_list)
    else:
        _build_category_chart(prs, option, series_list, chart_type)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
