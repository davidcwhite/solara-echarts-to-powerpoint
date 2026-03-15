"""Round-trip tests: ECharts option -> PPTX -> read back and verify."""

from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from converter import echarts_to_pptx, is_exportable
from theme import CORPORATE, DEFAULT, MONOCHROME, VIBRANT, PptxTheme

TEST_OUTPUT = Path(__file__).resolve().parent.parent / "test_output"


@pytest.fixture(autouse=True)
def _output_dir():
    TEST_OUTPUT.mkdir(exist_ok=True)


def _load_chart(pptx_bytes: bytes):
    """Load PPTX bytes and return the first chart object found."""
    prs = Presentation(BytesIO(pptx_bytes))
    assert len(prs.slides) == 1
    for shape in prs.slides[0].shapes:
        if shape.has_chart:
            return shape.chart
    pytest.fail("No chart shape found on the slide")


# ── Bar chart ────────────────────────────────────────────────────────

BAR_OPTION = {
    "title": {"text": "Bar Test"},
    "xAxis": {"type": "category", "data": ["Mon", "Tue", "Wed"]},
    "yAxis": {"type": "value"},
    "series": [{"name": "Sales", "type": "bar", "data": [120, 200, 150]}],
}


def test_bar_chart_type():
    data = echarts_to_pptx(BAR_OPTION)
    (TEST_OUTPUT / "bar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_bar_chart_categories():
    chart = _load_chart(echarts_to_pptx(BAR_OPTION))
    categories = [str(c) for c in chart.plots[0].categories]
    assert categories == ["Mon", "Tue", "Wed"]


def test_bar_chart_values():
    chart = _load_chart(echarts_to_pptx(BAR_OPTION))
    assert list(chart.series[0].values) == [120.0, 200.0, 150.0]


# ── Multi-series bar ────────────────────────────────────────────────

MULTI_BAR_OPTION = {
    "title": {"text": "Multi-series Bar"},
    "xAxis": {"type": "category", "data": ["A", "B", "C"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "Series 1", "type": "bar", "data": [10, 20, 30]},
        {"name": "Series 2", "type": "bar", "data": [15, 25, 35]},
    ],
}


def test_multi_series_bar():
    data = echarts_to_pptx(MULTI_BAR_OPTION)
    (TEST_OUTPUT / "multi_bar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert len(chart.series) == 2
    assert list(chart.series[0].values) == [10.0, 20.0, 30.0]
    assert list(chart.series[1].values) == [15.0, 25.0, 35.0]


# ── Bar with named data points ──────────────────────────────────────

NAMED_BAR_OPTION = {
    "title": {"text": "Named Bar"},
    "xAxis": {"type": "category"},
    "yAxis": {"type": "value"},
    "series": [
        {
            "name": "sales",
            "type": "bar",
            "data": [
                {"name": "Shirts", "value": 5},
                {"name": "Pants", "value": 10},
                {"name": "Socks", "value": 20},
            ],
        }
    ],
}


def test_named_data_categories():
    data = echarts_to_pptx(NAMED_BAR_OPTION)
    (TEST_OUTPUT / "named_bar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    categories = [str(c) for c in chart.plots[0].categories]
    assert categories == ["Shirts", "Pants", "Socks"]


def test_named_data_values():
    chart = _load_chart(echarts_to_pptx(NAMED_BAR_OPTION))
    assert list(chart.series[0].values) == [5.0, 10.0, 20.0]


# ── Stacked bar ─────────────────────────────────────────────────────

STACKED_BAR_OPTION = {
    "title": {"text": "Stacked Bar Test"},
    "xAxis": {"type": "category", "data": ["Q1", "Q2", "Q3"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "Online", "type": "bar", "stack": "total", "data": [100, 200, 300]},
        {"name": "In-Store", "type": "bar", "stack": "total", "data": [50, 80, 120]},
    ],
}


def test_stacked_bar_type():
    data = echarts_to_pptx(STACKED_BAR_OPTION)
    (TEST_OUTPUT / "stacked_bar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_STACKED


def test_stacked_bar_series_count():
    chart = _load_chart(echarts_to_pptx(STACKED_BAR_OPTION))
    assert len(chart.series) == 2


def test_stacked_bar_values():
    chart = _load_chart(echarts_to_pptx(STACKED_BAR_OPTION))
    assert list(chart.series[0].values) == [100.0, 200.0, 300.0]
    assert list(chart.series[1].values) == [50.0, 80.0, 120.0]


# ── Horizontal bar ──────────────────────────────────────────────────

HORIZONTAL_BAR_OPTION = {
    "title": {"text": "Horizontal Bar Test"},
    "xAxis": {"type": "value"},
    "yAxis": {"type": "category", "data": ["Speed", "UX", "Support"]},
    "series": [{"name": "Score", "type": "bar", "data": [85, 78, 88]}],
}


def test_horizontal_bar_type():
    data = echarts_to_pptx(HORIZONTAL_BAR_OPTION)
    (TEST_OUTPUT / "horizontal_bar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED


def test_horizontal_bar_categories():
    chart = _load_chart(echarts_to_pptx(HORIZONTAL_BAR_OPTION))
    categories = [str(c) for c in chart.plots[0].categories]
    assert categories == ["Speed", "UX", "Support"]


def test_horizontal_bar_values():
    chart = _load_chart(echarts_to_pptx(HORIZONTAL_BAR_OPTION))
    assert list(chart.series[0].values) == [85.0, 78.0, 88.0]


# ── Stacked horizontal bar ──────────────────────────────────────────

STACKED_HORIZONTAL_BAR_OPTION = {
    "title": {"text": "Stacked Horizontal Bar"},
    "xAxis": {"type": "value"},
    "yAxis": {"type": "category", "data": ["A", "B"]},
    "series": [
        {"name": "S1", "type": "bar", "stack": "t", "data": [10, 20]},
        {"name": "S2", "type": "bar", "stack": "t", "data": [30, 40]},
    ],
}


def test_stacked_horizontal_bar_type():
    data = echarts_to_pptx(STACKED_HORIZONTAL_BAR_OPTION)
    (TEST_OUTPUT / "stacked_horizontal_bar.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.BAR_STACKED


# ── Line chart ──────────────────────────────────────────────────────

LINE_OPTION = {
    "title": {"text": "Line Test"},
    "xAxis": {"type": "category", "data": ["Jan", "Feb", "Mar"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "Revenue", "type": "line", "data": [300, 450, 500]},
    ],
}


def test_line_chart_type():
    data = echarts_to_pptx(LINE_OPTION)
    (TEST_OUTPUT / "line_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.LINE_MARKERS


def test_line_chart_categories():
    chart = _load_chart(echarts_to_pptx(LINE_OPTION))
    assert [str(c) for c in chart.plots[0].categories] == ["Jan", "Feb", "Mar"]


def test_line_chart_values():
    chart = _load_chart(echarts_to_pptx(LINE_OPTION))
    assert list(chart.series[0].values) == [300.0, 450.0, 500.0]


# ── Area chart (line + areaStyle) ───────────────────────────────────

AREA_OPTION = {
    "title": {"text": "Area Test"},
    "xAxis": {"type": "category", "data": ["Mon", "Tue", "Wed"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "Traffic", "type": "line", "areaStyle": {}, "data": [820, 932, 901]},
    ],
}


def test_area_chart_type():
    data = echarts_to_pptx(AREA_OPTION)
    (TEST_OUTPUT / "area_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.AREA


def test_area_chart_values():
    chart = _load_chart(echarts_to_pptx(AREA_OPTION))
    assert list(chart.series[0].values) == [820.0, 932.0, 901.0]


# ── Stacked area ───────────────────────────────────────────────────

STACKED_AREA_OPTION = {
    "title": {"text": "Stacked Area"},
    "xAxis": {"type": "category", "data": ["A", "B"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "S1", "type": "line", "areaStyle": {}, "stack": "t", "data": [10, 20]},
        {"name": "S2", "type": "line", "areaStyle": {}, "stack": "t", "data": [30, 40]},
    ],
}


def test_stacked_area_type():
    data = echarts_to_pptx(STACKED_AREA_OPTION)
    (TEST_OUTPUT / "stacked_area_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.AREA_STACKED


# ── Pie chart ───────────────────────────────────────────────────────

PIE_OPTION = {
    "title": {"text": "Pie Test"},
    "series": [
        {
            "name": "Share",
            "type": "pie",
            "data": [
                {"name": "A", "value": 100},
                {"name": "B", "value": 200},
                {"name": "C", "value": 300},
            ],
        }
    ],
}


def test_pie_chart_type():
    data = echarts_to_pptx(PIE_OPTION)
    (TEST_OUTPUT / "pie_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.PIE


def test_pie_chart_categories():
    chart = _load_chart(echarts_to_pptx(PIE_OPTION))
    assert [str(c) for c in chart.plots[0].categories] == ["A", "B", "C"]


def test_pie_chart_values():
    chart = _load_chart(echarts_to_pptx(PIE_OPTION))
    assert list(chart.series[0].values) == [100.0, 200.0, 300.0]


# ── Donut chart (pie with inner radius) ─────────────────────────────

DONUT_OPTION = {
    "title": {"text": "Donut Test"},
    "series": [
        {
            "name": "Browsers",
            "type": "pie",
            "radius": ["40%", "70%"],
            "data": [
                {"name": "Chrome", "value": 65},
                {"name": "Firefox", "value": 12},
                {"name": "Safari", "value": 15},
            ],
        }
    ],
}


def test_donut_chart_type():
    data = echarts_to_pptx(DONUT_OPTION)
    (TEST_OUTPUT / "donut_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.DOUGHNUT


def test_donut_categories():
    chart = _load_chart(echarts_to_pptx(DONUT_OPTION))
    assert [str(c) for c in chart.plots[0].categories] == [
        "Chrome",
        "Firefox",
        "Safari",
    ]


def test_donut_values():
    chart = _load_chart(echarts_to_pptx(DONUT_OPTION))
    assert list(chart.series[0].values) == [65.0, 12.0, 15.0]


def test_pie_without_radius_is_not_donut():
    chart = _load_chart(echarts_to_pptx(PIE_OPTION))
    assert chart.chart_type == XL_CHART_TYPE.PIE


def test_pie_with_zero_inner_radius_is_not_donut():
    option = {
        "series": [
            {
                "type": "pie",
                "radius": ["0%", "60%"],
                "data": [{"name": "X", "value": 1}],
            }
        ],
    }
    chart = _load_chart(echarts_to_pptx(option))
    assert chart.chart_type == XL_CHART_TYPE.PIE


# ── Scatter chart ───────────────────────────────────────────────────

SCATTER_OPTION = {
    "title": {"text": "Scatter Test"},
    "xAxis": {"type": "value"},
    "yAxis": {"type": "value"},
    "series": [
        {
            "name": "Points",
            "type": "scatter",
            "data": [[1, 10], [2, 20], [3, 30]],
        }
    ],
}


def test_scatter_chart_type():
    data = echarts_to_pptx(SCATTER_OPTION)
    (TEST_OUTPUT / "scatter_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER


def test_scatter_series_count():
    chart = _load_chart(echarts_to_pptx(SCATTER_OPTION))
    assert len(chart.series) == 1


# ── effectScatter (treated as scatter) ──────────────────────────────

EFFECT_SCATTER_OPTION = {
    "title": {"text": "Effect Scatter"},
    "xAxis": {"type": "value"},
    "yAxis": {"type": "value"},
    "series": [
        {
            "name": "Ripples",
            "type": "effectScatter",
            "data": [[5, 50], [10, 100]],
        }
    ],
}


def test_effect_scatter_type():
    data = echarts_to_pptx(EFFECT_SCATTER_OPTION)
    (TEST_OUTPUT / "effect_scatter_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER


# ── Radar chart ─────────────────────────────────────────────────────

RADAR_OPTION = {
    "title": {"text": "Radar Test"},
    "radar": {
        "indicator": [
            {"name": "Eng", "max": 100},
            {"name": "Design", "max": 100},
            {"name": "Sales", "max": 100},
        ]
    },
    "series": [
        {
            "type": "radar",
            "data": [
                {"value": [90, 65, 80], "name": "Team A"},
                {"value": [70, 85, 60], "name": "Team B"},
            ],
        }
    ],
}


def test_radar_chart_type():
    data = echarts_to_pptx(RADAR_OPTION)
    (TEST_OUTPUT / "radar_chart.pptx").write_bytes(data)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.RADAR_MARKERS


def test_radar_categories():
    chart = _load_chart(echarts_to_pptx(RADAR_OPTION))
    categories = [str(c) for c in chart.plots[0].categories]
    assert categories == ["Eng", "Design", "Sales"]


def test_radar_series_count():
    chart = _load_chart(echarts_to_pptx(RADAR_OPTION))
    assert len(chart.series) == 2


def test_radar_values():
    chart = _load_chart(echarts_to_pptx(RADAR_OPTION))
    assert list(chart.series[0].values) == [90.0, 65.0, 80.0]
    assert list(chart.series[1].values) == [70.0, 85.0, 60.0]


# ── is_exportable ───────────────────────────────────────────────────


def test_exportable_bar():
    assert is_exportable(BAR_OPTION) is True


def test_exportable_line():
    assert is_exportable(LINE_OPTION) is True


def test_exportable_pie():
    assert is_exportable(PIE_OPTION) is True


def test_exportable_scatter():
    assert is_exportable(SCATTER_OPTION) is True


def test_exportable_effect_scatter():
    assert is_exportable(EFFECT_SCATTER_OPTION) is True


def test_exportable_radar():
    assert is_exportable(RADAR_OPTION) is True


def test_not_exportable_sunburst():
    option = {"series": [{"type": "sunburst", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_treemap():
    option = {"series": [{"type": "treemap", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_funnel():
    option = {"series": [{"type": "funnel", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_gauge():
    option = {"series": [{"type": "gauge", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_sankey():
    option = {"series": [{"type": "sankey", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_heatmap():
    option = {"series": [{"type": "heatmap", "data": []}]}
    assert is_exportable(option) is False


def test_not_exportable_empty():
    assert is_exportable({}) is False
    assert is_exportable({"series": []}) is False


# ── Edge cases ──────────────────────────────────────────────────────


def test_empty_series_raises():
    with pytest.raises(ValueError, match="at least one series"):
        echarts_to_pptx({"series": []})


def test_no_series_key_raises():
    with pytest.raises(ValueError, match="at least one series"):
        echarts_to_pptx({})


def test_no_title_does_not_crash():
    option = {
        "xAxis": {"type": "category", "data": ["X"]},
        "yAxis": {"type": "value"},
        "series": [{"type": "bar", "data": [42]}],
    }
    data = echarts_to_pptx(option)
    chart = _load_chart(data)
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_xaxis_as_list():
    """ECharts allows xAxis to be an array of axis objects."""
    option = {
        "xAxis": [{"type": "category", "data": ["Q1", "Q2"]}],
        "yAxis": {"type": "value"},
        "series": [{"type": "bar", "data": [10, 20]}],
    }
    chart = _load_chart(echarts_to_pptx(option))
    assert [str(c) for c in chart.plots[0].categories] == ["Q1", "Q2"]


# ── Theme tests ─────────────────────────────────────────────────────

THEMED_BAR = {
    "title": {"text": "Themed Bar"},
    "xAxis": {"type": "category", "data": ["A", "B"]},
    "yAxis": {"type": "value"},
    "series": [
        {"name": "S1", "type": "bar", "data": [10, 20]},
        {"name": "S2", "type": "bar", "data": [30, 40]},
    ],
}


def _load_prs(pptx_bytes: bytes):
    return Presentation(BytesIO(pptx_bytes))


def test_default_theme_bar_series_colors():
    data = echarts_to_pptx(THEMED_BAR, theme=DEFAULT)
    (TEST_OUTPUT / "themed_default_bar.pptx").write_bytes(data)
    chart = _load_chart(data)
    c0 = chart.series[0].format.fill.fore_color.rgb
    c1 = chart.series[1].format.fill.fore_color.rgb
    assert c0 == RGBColor.from_string(DEFAULT.palette[0])
    assert c1 == RGBColor.from_string(DEFAULT.palette[1])


def test_corporate_theme_bar_series_colors():
    data = echarts_to_pptx(THEMED_BAR, theme=CORPORATE)
    (TEST_OUTPUT / "themed_corporate_bar.pptx").write_bytes(data)
    chart = _load_chart(data)
    c0 = chart.series[0].format.fill.fore_color.rgb
    assert c0 == RGBColor.from_string(CORPORATE.palette[0])


def test_vibrant_theme_bar_series_colors():
    data = echarts_to_pptx(THEMED_BAR, theme=VIBRANT)
    (TEST_OUTPUT / "themed_vibrant_bar.pptx").write_bytes(data)
    chart = _load_chart(data)
    c0 = chart.series[0].format.fill.fore_color.rgb
    assert c0 == RGBColor.from_string(VIBRANT.palette[0])


def test_monochrome_theme_bar_series_colors():
    data = echarts_to_pptx(THEMED_BAR, theme=MONOCHROME)
    (TEST_OUTPUT / "themed_monochrome_bar.pptx").write_bytes(data)
    chart = _load_chart(data)
    c0 = chart.series[0].format.fill.fore_color.rgb
    assert c0 == RGBColor.from_string(MONOCHROME.palette[0])


def test_theme_title_font():
    data = echarts_to_pptx(THEMED_BAR, theme=CORPORATE)
    prs = _load_prs(data)
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            p = shape.text_frame.paragraphs[0]
            assert p.font.name == CORPORATE.title_font
            assert p.font.size.pt == CORPORATE.title_size
            break
    else:
        pytest.fail("No title textbox found")


def test_theme_legend_position_bottom():
    data = echarts_to_pptx(THEMED_BAR, theme=DEFAULT)
    chart = _load_chart(data)
    assert chart.has_legend
    assert chart.legend.position == XL_LEGEND_POSITION.BOTTOM


def test_theme_legend_position_right():
    data = echarts_to_pptx(THEMED_BAR, theme=CORPORATE)
    chart = _load_chart(data)
    assert chart.has_legend
    assert chart.legend.position == XL_LEGEND_POSITION.RIGHT


def test_line_chart_uses_line_color():
    option = {
        "title": {"text": "Themed Line"},
        "xAxis": {"type": "category", "data": ["A", "B"]},
        "yAxis": {"type": "value"},
        "series": [{"name": "L1", "type": "line", "data": [10, 20]}],
    }
    data = echarts_to_pptx(option, theme=VIBRANT)
    (TEST_OUTPUT / "themed_vibrant_line.pptx").write_bytes(data)
    chart = _load_chart(data)
    line_color = chart.series[0].format.line.color.rgb
    assert line_color == RGBColor.from_string(VIBRANT.palette[0])


def test_pie_chart_point_colors():
    option = {
        "series": [
            {
                "type": "pie",
                "data": [
                    {"name": "X", "value": 10},
                    {"name": "Y", "value": 20},
                ],
            }
        ],
    }
    data = echarts_to_pptx(option, theme=DEFAULT)
    (TEST_OUTPUT / "themed_default_pie.pptx").write_bytes(data)
    chart = _load_chart(data)
    p0 = chart.plots[0].series[0].points[0].format.fill.fore_color.rgb
    p1 = chart.plots[0].series[0].points[1].format.fill.fore_color.rgb
    assert p0 == RGBColor.from_string(DEFAULT.palette[0])
    assert p1 == RGBColor.from_string(DEFAULT.palette[1])


def test_all_themes_produce_valid_pptx():
    """Smoke test: every built-in theme produces a loadable PPTX."""
    for theme in (DEFAULT, CORPORATE, VIBRANT, MONOCHROME):
        data = echarts_to_pptx(BAR_OPTION, theme=theme)
        prs = _load_prs(data)
        assert len(prs.slides) == 1
